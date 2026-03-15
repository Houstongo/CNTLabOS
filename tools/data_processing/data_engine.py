import os
import re
import zipfile
import xml.etree.ElementTree as ET
import pandas as pd
from pptx import Presentation
from typing import Dict, Any, List

# 项目路径配置 (遵循 context.md 规则)
BASE_DIR = r"d:\CNTDATA\CNTA_ML_Project"
TEMP_DIR = os.path.join(BASE_DIR, "data", "temp")
os.makedirs(TEMP_DIR, exist_ok=True)

class DataEngine:
    """
    集成的实验数据处理引擎：支持文件夹名解析、DOCX/PPTX提取、Excel处理。
    """
    
    @staticmethod
    def parse_exp_folder(folder_name: str) -> Dict[str, Any]:
        """解析文件夹名中的实验参数"""
        params = {
            'folder_name': folder_name,
            'date': None, 'temp_C': None, 'time_h': None, 'flow_L': None, 'weight_g': None
        }
        # 日期 (前面6位)
        date_m = re.search(r'^(\d{6})', folder_name)
        if date_m: params['date'] = date_m.group(1)
        
        # 温度 (Txxx)
        temp_m = re.search(r'[tT](\d+)', folder_name)
        if temp_m: params['temp_C'] = int(temp_m.group(1))
        
        # 时间 (xh)
        time_m = re.search(r'(\d+)h', folder_name, re.IGNORECASE)
        if time_m: params['time_h'] = int(time_m.group(1))
        
        # 流量 (Lxxx)
        flow_m = re.search(r'[lL](\d+)', folder_name)
        if flow_m: params['flow_L'] = int(flow_m.group(1))
        
        # 质量 (x.xg)
        weight_m = re.search(r'(\d+\.?\d*)g', folder_name)
        if weight_m: params['weight_g'] = float(weight_m.group(1))
        
        return params

    @staticmethod
    def extract_docx_text(file_path: str, output_name: str = "doc_content.txt"):
        """提取 Word 文档文本"""
        print(f"解析 DOCX: {os.path.basename(file_path)}")
        try:
            with zipfile.ZipFile(file_path) as docx:
                xml_content = docx.read('word/document.xml')
                tree = ET.fromstring(xml_content)
                ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                paragraphs = [''.join(t.text for t in p.findall('.//w:t', ns) if t.text) 
                             for p in tree.findall('.//w:p', ns)]
                text = '\n'.join(filter(None, paragraphs))
                
                out_path = os.path.join(TEMP_DIR, output_name)
                with open(out_path, "w", encoding="utf-8") as f:
                    f.write(text)
                print(f"成功导出文本至: {out_path}")
        except Exception as e:
            print(f"DOCX 处理失败: {e}")

    @staticmethod
    def process_excel(file_path: str):
        """处理 Excel 数据并保存 CSV 副本"""
        print(f"处理 Excel: {os.path.basename(file_path)}")
        try:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in df_dict.items():
                out_path = os.path.join(TEMP_DIR, f"{sheet_name}.csv")
                df.to_csv(out_path, index=False, encoding='utf-8-sig')
                print(f"  - 已将 sheet [{sheet_name}] 转换为 {out_path}")
        except Exception as e:
            print(f"Excel 处理失败: {e}")

    @staticmethod
    def process_pptx(file_path: str):
        """提取 PPTX 文本内容"""
        print(f"处理 PPTX: {os.path.basename(file_path)}")
        try:
            prs = Presentation(file_path)
            content = []
            for i, slide in enumerate(prs.slides):
                slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text")]
                if slide_text:
                    content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
            
            out_path = os.path.join(TEMP_DIR, os.path.basename(file_path) + ".txt")
            with open(out_path, "w", encoding="utf-8") as f:
                f.write("\n\n".join(content))
            print(f"已导出 PPTX 内容至: {out_path}")
        except Exception as e:
            print(f"PPTX 处理失败: {e}")

def run_full_etl():
    """主 ETL 流程"""
    engine = DataEngine()
    
    # 示例外部路径 (可根据需要调整)
    ext_files = [
        r"d:\CNTDATA\doc\3机器学习辅助的碳纳米管阵列的数字特征提取及生长工艺优化(1).docx",
        r"D:\电脑管家迁移文件\xwechat_files\wxid_697jsgpwq0zg22_414a\msg\file\2026-03\管式炉温区.xlsx",
        r"D:\电脑管家迁移文件\xwechat_files\wxid_697jsgpwq0zg22_414a\msg\file\2026-03\碳管温区.pptx"
    ]

    print("=== CNTA 项目数据集成引擎启动 ===")
    for f in ext_files:
        if not os.path.exists(f): 
            print(f"跳过不存在的文件: {f}")
            continue
        
        if f.endswith('.docx'): engine.extract_docx_text(f)
        elif f.endswith('.xlsx'): engine.process_excel(f)
        elif f.endswith('.pptx'): engine.process_pptx(f)

    print("\n=== 实验批次索引生成 ===")
    root_dir = r"d:\CNTDATA"
    all_exps = []
    for item in os.listdir(root_dir):
        if os.path.isdir(os.path.join(root_dir, item)) and item != 'CNTA_ML_Project':
            all_exps.append(engine.parse_exp_folder(item))
    
    if all_exps:
        df = pd.DataFrame(all_exps)
        idx_path = os.path.join(BASE_DIR, "data", "dataset_index.csv")
        df.to_csv(idx_path, index=False, encoding='utf-8-sig')
        print(f"索引已更新: {idx_path}")

if __name__ == "__main__":
    run_full_etl()
