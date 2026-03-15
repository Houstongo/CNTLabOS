import pandas as pd
from pptx import Presentation
import os

files = [
    r"D:\电脑管家迁移文件\xwechat_files\wxid_697jsgpwq0zg22_414a\msg\file\2026-03\管式炉温区.xlsx",
    r"D:\电脑管家迁移文件\xwechat_files\wxid_697jsgpwq0zg22_414a\msg\file\2026-03\碳管温区.pptx"
]

def read_xlsx(path):
    print(f"\n--- Reading Excel: {os.path.basename(path)} ---")
    try:
        df = pd.read_excel(path, sheet_name=None)
        for sheet_name, data in df.items():
            print(f"\nSheet: {sheet_name}")
            print(data.head())
    except Exception as e:
        print(f"Error reading Excel: {e}")

def read_pptx(path):
    print(f"\n--- Reading PPTX: {os.path.basename(path)} ---")
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            print(f"\nSlide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text.strip())
    except Exception as e:
        print(f"Error reading PPTX: {e}")

for file in files:
    if file.endswith('.xlsx'):
        read_xlsx(file)
    elif file.endswith('.pptx'):
        read_pptx(file)
